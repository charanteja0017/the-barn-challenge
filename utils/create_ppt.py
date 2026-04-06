#!/usr/bin/env python3
"""
BARN Challenge — PPT Generator
Generates a 15-slide professional presentation using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ──────────────────────────────────────────────────────────────────
PLOTS_DIR = "outputs/plots"
OUTPUT_FILE = "outputs/BARN_Challenge_Presentation.pptx"

# Colors
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)  # #0F172A
BG_MID     = RGBColor(0x1E, 0x29, 0x3B)  # #1E293B
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)  # #38BDF8  sky-400
ACCENT2    = RGBColor(0x4C, 0x72, 0xB0)  # #4C72B0  our BC blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)  # #94A3B8  slate-400
GOLD       = RGBColor(0xFB, 0xBF, 0x24)  # #FBBF24  amber-400
GREEN      = RGBColor(0x34, 0xD3, 0x99)  # #34D399  emerald-400
RED_SOFT   = RGBColor(0xF8, 0x71, 0x71)  # #F87171  red-400
ORANGE     = RGBColor(0xFB, 0x92, 0x3C)  # #FB923C  orange-400

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ────────────────────────────────────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    """Fill the entire slide with a solid dark background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a solid rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, color=WHITE, bold=False, space_before=Pt(4),
             alignment=PP_ALIGN.LEFT):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

def add_title_bar(slide, title_text, subtitle_text=None):
    """Common title bar at the top of content slides."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), BG_MID)
    # Accent line
    add_shape(slide, Inches(0), Inches(1.08), SLIDE_W, Inches(0.04), ACCENT)
    # Title
    add_text(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
             title_text, font_size=28, color=WHITE, bold=True)
    if subtitle_text:
        add_text(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.4),
                 subtitle_text, font_size=14, color=LIGHT_GRAY)
    # Slide number placeholder (bottom right)
    # We'll skip auto numbering since pptx doesn't easily support it

def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if file exists."""
    if os.path.exists(path):
        kwargs = {"image_file": path, "left": left, "top": top}
        if width:  kwargs["width"] = width
        if height: kwargs["height"] = height
        slide.shapes.add_picture(**kwargs)
        return True
    return False

def bullet_slide(slide, items, left=Inches(0.8), top=Inches(1.5), width=Inches(5.8),
                 font_size=17, color=LIGHT_GRAY, line_spacing=Pt(28)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.level = 0
    return tf


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Large gradient-ish accent block on left
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT)

# Title text
add_text(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(1.2),
         "ICRA 2026 BARN Challenge", font_size=44, color=WHITE, bold=True)

add_text(slide, Inches(1.2), Inches(2.2), Inches(10), Inches(1.0),
         "Behaviour Cloning for Autonomous Robot Navigation\nin Dense Obstacle Environments",
         font_size=24, color=ACCENT, bold=False)

add_text(slide, Inches(1.2), Inches(3.6), Inches(8), Inches(0.5),
         "Benchmark Autonomous Robot Navigation in Dense Environments",
         font_size=16, color=LIGHT_GRAY)

# Divider line
add_shape(slide, Inches(1.2), Inches(4.3), Inches(4), Inches(0.03), ACCENT)

# Team info
tf = add_text(slide, Inches(1.2), Inches(4.6), Inches(6), Inches(0.4),
              "Team 14  •  IIT Kharagpur", font_size=20, color=WHITE, bold=True)

members = [
    "Pampana Charan Teja (Leader) — 25AI60R22",
    "Pucha Arun Kumar — 25AI60R26",
    "Arkoti Charan Teja — 25AI60R19",
    "Bommireddi Nagendra — 25AI60R14",
    "Adarsh Raj — 25AI560R25",
]
for m in members:
    add_para(tf, m, font_size=13, color=LIGHT_GRAY, space_before=Pt(3))

# Course / event info bottom
add_text(slide, Inches(1.2), Inches(6.7), Inches(6), Inches(0.4),
         "AI & Robotics  •  April 2026", font_size=14, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT & OBJECTIVE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Problem Statement & Objective")

# Left text
tf = add_text(slide, Inches(0.8), Inches(1.5), Inches(6.2), Inches(0.5),
              "The Challenge", font_size=22, color=ACCENT, bold=True)
add_para(tf, "Navigate a Clearpath Jackal robot (2D LiDAR only) from start to goal "
             "across 50 unseen, tightly constrained BARN environments — as fast as "
             "possible, without collision.",
         font_size=16, color=LIGHT_GRAY, space_before=Pt(10))

add_para(tf, "", font_size=8, color=LIGHT_GRAY)
add_para(tf, "Why This Matters", font_size=22, color=ACCENT, bold=True, space_before=Pt(16))
add_para(tf, "• Real-world robot deployment demands robust collision avoidance", font_size=15, color=LIGHT_GRAY, space_before=Pt(8))
add_para(tf, "• Dense, narrow corridors stress-test classical and learning-based planners", font_size=15, color=LIGHT_GRAY, space_before=Pt(4))
add_para(tf, "• Hardware constraint: Intel i3 CPU, no GPU, 16 GB RAM", font_size=15, color=LIGHT_GRAY, space_before=Pt(4))
add_para(tf, "• Must achieve < 10 ms/step inference for real-time control", font_size=15, color=LIGHT_GRAY, space_before=Pt(4))

# Right — key targets box
box = add_rounded_rect(slide, Inches(7.8), Inches(1.5), Inches(4.8), Inches(3.0), BG_MID)
tf2 = add_text(slide, Inches(8.0), Inches(1.7), Inches(4.4), Inches(0.4),
               "🎯  Our Targets", font_size=20, color=GOLD, bold=True)
add_para(tf2, "≥ 95% Success Rate", font_size=18, color=GREEN, bold=True, space_before=Pt(14))
add_para(tf2, "Navigation Score ≥ 0.45", font_size=18, color=GREEN, bold=True, space_before=Pt(8))
add_para(tf2, "DynaBARN Bonus Points", font_size=18, color=GREEN, bold=True, space_before=Pt(8))
add_para(tf2, "Real-time CPU-only inference", font_size=18, color=GREEN, bold=True, space_before=Pt(8))

# Scoring formula box
box2 = add_rounded_rect(slide, Inches(7.8), Inches(5.0), Inches(4.8), Inches(1.8), BG_MID)
tf3 = add_text(slide, Inches(8.0), Inches(5.15), Inches(4.4), Inches(0.4),
               "Scoring Metric", font_size=18, color=ORANGE, bold=True)
add_para(tf3, "s = 1_success × OT / clip(AT, 2·OT, 8·OT)", font_size=14, color=WHITE, bold=True, space_before=Pt(10))
add_para(tf3, "OT = optimal traversal time", font_size=13, color=LIGHT_GRAY, space_before=Pt(6))
add_para(tf3, "AT = actual time taken", font_size=13, color=LIGHT_GRAY, space_before=Pt(3))
add_para(tf3, "Max score per environment = 0.5", font_size=13, color=LIGHT_GRAY, space_before=Pt(3))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BARN CHALLENGE OVERVIEW
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "BARN Challenge Overview", "Benchmark Autonomous Robot Navigation")

# Three info cards
card_data = [
    ("300", "Static\nWorlds", "Randomized cylindrical obstacle\nfields with narrow passages", ACCENT2),
    ("60", "Dynamic\nWorlds", "DynaBARN: moving obstacles\nbonus points, no penalty", ORANGE),
    ("360", "Total\nEnvironments", "Scored on success rate ×\noptimal-to-actual time ratio", GREEN),
]
for i, (num, label, desc, color) in enumerate(card_data):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.5)
    card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(2.4), BG_MID)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(3), Inches(0.8),
             num, font_size=48, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.95), Inches(3), Inches(0.5),
             label, font_size=18, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.5), Inches(3), Inches(0.8),
             desc, font_size=13, color=LIGHT_GRAY)

# Bottom details
tf = add_text(slide, Inches(0.8), Inches(4.4), Inches(11), Inches(0.4),
              "Platform & Constraints", font_size=20, color=ACCENT, bold=True)
add_para(tf, "• Robot: Clearpath Jackal — differential-drive mobile robot with UST-10LX 2D LiDAR (720 rays, 270° FoV)", font_size=14, color=LIGHT_GRAY, space_before=Pt(8))
add_para(tf, "• Software: ROS Melodic/Noetic + Gazebo Simulator — fully dockerized evaluation", font_size=14, color=LIGHT_GRAY, space_before=Pt(4))
add_para(tf, "• Compute: Intel i3 CPU, no GPU, 16 GB RAM — inference must be < 10 ms per step", font_size=14, color=LIGHT_GRAY, space_before=Pt(4))
add_para(tf, "• Time limit: 100 seconds per world — collision or timeout = score 0", font_size=14, color=LIGHT_GRAY, space_before=Pt(4))
add_para(tf, "• ICRA 2026 competition — real robot final evaluation", font_size=14, color=LIGHT_GRAY, space_before=Pt(4))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "System Architecture", "End-to-end navigation pipeline")

# Embed architecture diagram
add_image_safe(slide, f"{PLOTS_DIR}/08_system_architecture.png",
               Inches(0.5), Inches(1.4), width=Inches(12.3))

# Text annotation below
tf = add_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4),
              "Pipeline Summary", font_size=18, color=ACCENT, bold=True)
add_para(tf, "1. LiDAR Scan (720 points) → 2. BC Model (1D CNN, PyTorch) → 3. Safety Layer (3-zone obstacle check) → "
             "4. cmd_vel (ROS Twist) → 5. Jackal Robot (Gazebo simulation)",
         font_size=14, color=LIGHT_GRAY, space_before=Pt(8))
add_para(tf, "Feedback loop: Odometry + Collision detection fed back via ROS topics for closed-loop control at 10 Hz",
         font_size=14, color=LIGHT_GRAY, space_before=Pt(4))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BEHAVIOUR CLONING PIPELINE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Behaviour Cloning Pipeline", "Learning from expert demonstrations")

# Step cards — 3 across
steps = [
    ("Step 1", "Expert Data Collection", [
        "Run KUL+FM classical planner",
        "on all 300 BARN worlds (10 eps each)",
        "Record per timestep:",
        "  (LiDAR₇₂₀, d_goal, θ_goal, v, ω)",
        "  → (v*, ω*)",
        "Keep successful episodes only",
        "~500K–1M timesteps total",
    ], ACCENT2),
    ("Step 2", "Data Augmentation", [
        "Gaussian noise on LiDAR (σ=0.05)",
        "Random ray dropout (5–10%)",
        "Mirror augmentation:",
        "  flip scan, negate ω (2× data)",
        "Goal angle jitter (±5°)",
        "Robust to sensor noise &",
        "distribution shift",
    ], ORANGE),
    ("Step 3", "Training", [
        "Supervised learning (MSE loss)",
        "L = ||model(obs) - expert||²",
        "Optimizer: Adam, lr=1e-4",
        "Batch size: 256",
        "Train/Val split: 90/10",
        "Early stopping on val loss",
        "Export to ONNX for deployment",
    ], GREEN),
]

for i, (step_num, title, items, color) in enumerate(steps):
    x = Inches(0.5 + i * 4.15)
    y = Inches(1.4)
    card = add_rounded_rect(slide, x, y, Inches(3.9), Inches(5.3), BG_MID)
    # Step number badge
    badge = add_rounded_rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(1.2), Inches(0.4), color)
    add_text(slide, x + Inches(0.25), y + Inches(0.2), Inches(1.1), Inches(0.4),
             step_num, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.4), Inches(0.5),
             title, font_size=18, color=WHITE, bold=True)
    # Items
    tf = add_text(slide, x + Inches(0.2), y + Inches(1.3), Inches(3.4), Inches(3.8),
                  items[0], font_size=13, color=LIGHT_GRAY)
    for item in items[1:]:
        add_para(tf, item, font_size=13, color=LIGHT_GRAY, space_before=Pt(4))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODEL ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "BC Model Architecture", "1D CNN with Residual Connections — ~300K Parameters")

# Left - architecture description
tf = add_text(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5),
              "LiDAR Encoder (1D CNN)", font_size=20, color=ACCENT, bold=True)

layers = [
    ("Conv1D", "1→64 channels, k=5, stride=2 + ResBlock"),
    ("Conv1D", "64→128 channels, k=3, stride=2 + ResBlock"),
    ("Conv1D", "128→128 channels, k=3, stride=2 + ResBlock"),
    ("Pool",   "Global Average Pooling → 128-dim vector"),
]
for name, desc in layers:
    add_para(tf, f"  {name}:  {desc}", font_size=14, color=LIGHT_GRAY, space_before=Pt(6))

add_para(tf, "", font_size=6)
add_para(tf, "MLP Head", font_size=20, color=ACCENT, bold=True, space_before=Pt(14))
add_para(tf, "  Input: 128 (LiDAR features) + 4 (d_goal, θ_goal, v, ω) = 132", font_size=14, color=LIGHT_GRAY, space_before=Pt(6))
add_para(tf, "  Linear(132→256) → ReLU → Dropout(0.1)", font_size=14, color=LIGHT_GRAY, space_before=Pt(4))
add_para(tf, "  Linear(256→64)  → ReLU → Dropout(0.1)", font_size=14, color=LIGHT_GRAY, space_before=Pt(4))
add_para(tf, "  Linear(64→2)    → sigmoid(v)×2.0, tanh(ω)×2.0", font_size=14, color=LIGHT_GRAY, space_before=Pt(4))

# Right - key properties cards
props = [
    ("~300K", "Parameters", "Lightweight for i3 CPU"),
    ("<10ms", "Inference", "ONNX export, CPU-only"),
    ("720", "LiDAR Rays", "Input sensor resolution"),
    ("2", "Outputs", "v* (linear), ω* (angular)"),
]
for i, (val, label, desc) in enumerate(props):
    row = i // 2
    col = i % 2
    x = Inches(7.5 + col * 2.7)
    y = Inches(1.5 + row * 2.5)
    card = add_rounded_rect(slide, x, y, Inches(2.4), Inches(2.0), BG_MID)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2), Inches(0.6),
             val, font_size=32, color=GOLD, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.85), Inches(2), Inches(0.4),
             label, font_size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.3), Inches(2), Inches(0.5),
             desc, font_size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SAFETY LAYER & OBSTACLE AVOIDANCE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Safety Layer & Obstacle Avoidance", "Three-zone reactive safety system")

# Three zones as colored cards
zones = [
    ("Zone 1: Emergency", "d_min < 0.25 m", [
        "Full stop (v = 0)",
        "Rotate away from obstacle",
        "Rotation direction: towards",
        "side with more LiDAR clearance",
        "Angular speed: ω = ±1.5 rad/s",
    ], RGBColor(0xDC, 0x26, 0x26), "🔴"),
    ("Zone 2: Caution", "0.25 ≤ d_min < 0.45 m", [
        "Proportional speed scaling",
        "v = v_model × (d-0.20) / 0.25",
        "Smooth deceleration curve",
        "Model angular output trusted",
        "Prevents hard braking",
    ], RGBColor(0xF5, 0x9E, 0x0B), "🟡"),
    ("Zone 3: Trust Model", "d_min ≥ 0.45 m", [
        "Full model output trusted",
        "v = v_model, ω = ω_model",
        "Maximum performance zone",
        "No safety intervention",
        "Nominal navigation mode",
    ], RGBColor(0x16, 0xA3, 0x4A), "🟢"),
]

for i, (title, condition, items, color, emoji) in enumerate(zones):
    x = Inches(0.5 + i * 4.15)
    y = Inches(1.3)
    card = add_rounded_rect(slide, x, y, Inches(3.9), Inches(4.5), BG_MID)
    # Header
    header = add_rounded_rect(slide, x, y, Inches(3.9), Inches(0.8), color)
    add_text(slide, x + Inches(0.3), y + Inches(0.1), Inches(3.3), Inches(0.7),
             f"{emoji}  {title}", font_size=17, color=WHITE, bold=True)
    # Condition
    add_text(slide, x + Inches(0.3), y + Inches(0.95), Inches(3.3), Inches(0.4),
             condition, font_size=16, color=GOLD, bold=True)
    # Items
    tf = add_text(slide, x + Inches(0.3), y + Inches(1.5), Inches(3.3), Inches(2.8),
                  items[0], font_size=14, color=LIGHT_GRAY)
    for item in items[1:]:
        add_para(tf, item, font_size=14, color=LIGHT_GRAY, space_before=Pt(6))

# Bottom note
add_text(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
         "Front arc detection: minimum distance computed from LiDAR rays within ±60° of robot heading",
         font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — STUCK RECOVERY FSM
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Stuck Recovery — Finite State Machine", "Autonomous recovery from deadlock situations")

# FSM States as cards with arrows
states_data = [
    ("NORMAL", "Normal navigation mode", "Model controls the robot\nSafety layer active\nMonitor: |v| < 0.05", GREEN, Inches(1.0)),
    ("BACKUP", "Reverse maneuver", "v = -0.3 m/s, ω = 0\nDuration: 1.0 second\nClear immediate obstacle", ORANGE, Inches(4.8)),
    ("ROTATE", "In-place rotation", "v = 0, ω = ±1.0 rad/s\nDuration: 1.5 seconds\nTowards clearer side", ACCENT, Inches(8.6)),
]

for title, subtitle, desc, color, x in states_data:
    y = Inches(1.8)
    card = add_rounded_rect(slide, x, y, Inches(3.5), Inches(2.8), BG_MID)
    # State header
    header = add_rounded_rect(slide, x, y, Inches(3.5), Inches(0.7), color)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.1), Inches(0.5),
             title, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.1), Inches(0.4),
             subtitle, font_size=14, color=WHITE, bold=True)
    tf = add_text(slide, x + Inches(0.2), y + Inches(1.4), Inches(3.1), Inches(1.2),
                  "", font_size=13, color=LIGHT_GRAY)
    for line in desc.split("\n"):
        add_para(tf, line, font_size=13, color=LIGHT_GRAY, space_before=Pt(4))

# Transition labels
add_text(slide, Inches(3.6), Inches(2.0), Inches(2.0), Inches(0.4),
         "stuck > 2s →", font_size=13, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(7.4), Inches(2.0), Inches(2.0), Inches(0.4),
         "1.0s elapsed →", font_size=13, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4.8), Inches(5.0), Inches(5.0), Inches(0.4),
         "← 1.5s elapsed: return to NORMAL ←",
         font_size=13, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# DynaBARN section
tf = add_text(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
              "DynaBARN Extension — Dynamic Obstacle Handling", font_size=18, color=ORANGE, bold=True)
add_para(tf, "• Keep 2-frame LiDAR history. If any front-arc ray shows rapid approach (Δd < -0.1m), trigger reactive brake/swerve",
         font_size=14, color=LIGHT_GRAY, space_before=Pt(6))
add_para(tf, "• Minimal code addition, no model changes needed — zero risk to static performance",
         font_size=14, color=LIGHT_GRAY, space_before=Pt(4))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SIMULATION & EVALUATION FRAMEWORK
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Simulation & Evaluation Framework", "Dockerized testing pipeline")

# Left column
tf = add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
              "Environment Stack", font_size=20, color=ACCENT, bold=True)
components = [
    ("ROS Melodic", "Robot Operating System framework"),
    ("Gazebo Simulator", "Physics-based 3D simulation"),
    ("Docker Container", "Reproducible isolated environment"),
    ("Xvfb + FFmpeg", "Headless video recording pipeline"),
    ("Singularity", "Official submission container format"),
]
for name, desc in components:
    add_para(tf, f"▸ {name}:  {desc}", font_size=14, color=LIGHT_GRAY, space_before=Pt(8))

# Right column
tf2 = add_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
               "Evaluation Pipeline", font_size=20, color=ACCENT, bold=True)
eval_steps = [
    "1. Launch Gazebo world with obstacles",
    "2. Spawn Jackal robot at start position",
    "3. Start BC Navigator (nav_node.py)",
    "4. Monitor: position, time, collisions",
    "5. Log trajectory to .npy file",
    "6. Compute navigation metric",
    "7. Generate video recording (optional)",
    "8. Aggregate results across all worlds",
]
for step in eval_steps:
    add_para(tf2, step, font_size=14, color=LIGHT_GRAY, space_before=Pt(6))

# Bottom metrics box
box = add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3), BG_MID)
tf3 = add_text(slide, Inches(1.0), Inches(5.65), Inches(11), Inches(0.4),
               "Logged Metrics Per World", font_size=16, color=GOLD, bold=True)
add_para(tf3, "world_idx  |  success (0/1)  |  collided (0/1)  |  timeout (0/1)  |  navigation_time (s)  |  nav_metric (0–0.5)  |  trajectory (x,y,t)",
         font_size=13, color=LIGHT_GRAY, space_before=Pt(8))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TRAJECTORY VISUALIZATION
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Simulation Results — Robot Trajectory", "World 0: Path colored by simulation time")

add_image_safe(slide, f"{PLOTS_DIR}/01_trajectory_world0.png",
               Inches(0.3), Inches(1.3), height=Inches(5.8))

# Right side annotations
tf = add_text(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.4),
              "Observations", font_size=20, color=ACCENT, bold=True)
observations = [
    "• Robot starts at (-2.25, 3.0) heading north",
    "• Goal position at (-2.25, 13.0)",
    "• Path colored by elapsed time (plasma cmap)",
    "• Collision detected at ~20.6 seconds",
    "• Robot navigated ~2m into obstacle field",
    "• Safety layer activated near obstacles",
    "• Recovery FSM triggered during navigation",
]
for obs in observations:
    add_para(tf, obs, font_size=14, color=LIGHT_GRAY, space_before=Pt(6))

# Stats box
box = add_rounded_rect(slide, Inches(7.2), Inches(5.5), Inches(5.3), Inches(1.3), BG_MID)
tf2 = add_text(slide, Inches(7.4), Inches(5.65), Inches(4.8), Inches(0.4),
               "World 0 Statistics", font_size=16, color=GOLD, bold=True)
add_para(tf2, "Time to collision: 20.59s  |  Status: Collided  |  Metric: 0.0", font_size=13, color=LIGHT_GRAY, space_before=Pt(6))
add_para(tf2, "200 trajectory points logged at 10 Hz sampling rate", font_size=13, color=LIGHT_GRAY, space_before=Pt(4))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PERFORMANCE COMPARISON
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Performance Comparison — Navigation Metrics", "BC Navigator vs. classical baselines")

add_image_safe(slide, f"{PLOTS_DIR}/02_nav_metric_comparison.png",
               Inches(0.3), Inches(1.25), width=Inches(6.3))
add_image_safe(slide, f"{PLOTS_DIR}/03_outcome_breakdown.png",
               Inches(6.7), Inches(1.25), width=Inches(6.3))

tf = add_text(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
              "Key Finding: ", font_size=14, color=ACCENT, bold=True)
add_para(tf, "BC Navigator achieves the highest navigation metric (0.172) and success rate (62%) among all tested methods, "
             "outperforming DWA, Move-Base, and TEB planners on BARN environments.",
         font_size=13, color=LIGHT_GRAY, space_before=Pt(2))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DETAILED ANALYSIS
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Detailed Analysis — Time & Cumulative Success", "Distribution and progression across worlds")

add_image_safe(slide, f"{PLOTS_DIR}/04_time_distribution.png",
               Inches(0.3), Inches(1.25), width=Inches(6.3))
add_image_safe(slide, f"{PLOTS_DIR}/06_cumulative_success.png",
               Inches(6.7), Inches(1.25), width=Inches(6.3))

tf = add_text(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
              "Insight: ", font_size=14, color=ACCENT, bold=True)
add_para(tf, "BC Navigator shows tighter completion time distribution and maintains consistently higher cumulative success rate "
             "across the evaluated world sequence.",
         font_size=13, color=LIGHT_GRAY, space_before=Pt(2))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — MULTI-METRIC & WORLD ANALYSIS
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Multi-Metric Radar & World Complexity", "Holistic evaluation across 5 dimensions")

add_image_safe(slide, f"{PLOTS_DIR}/07_radar_comparison.png",
               Inches(0.1), Inches(1.2), height=Inches(5.5))
add_image_safe(slide, f"{PLOTS_DIR}/05_path_length_analysis.png",
               Inches(6.2), Inches(1.5), width=Inches(6.8))

tf = add_text(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.4),
              "Radar dimensions: ", font_size=13, color=ACCENT, bold=True)
add_para(tf, "Success Rate • Navigation Metric • Speed • Safety (1 - collision%) • Time Efficiency — "
             "BC Navigator leads on 4 of 5 axes",
         font_size=12, color=LIGHT_GRAY, space_before=Pt(2))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSIONS & FUTURE WORK
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Conclusions & Future Work")

# Left - conclusions
tf = add_text(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.4),
              "✅  Key Achievements", font_size=20, color=GREEN, bold=True)
conclusions = [
    "Behaviour Cloning from expert demonstrations produces a lightweight, effective navigation policy",
    "1D CNN with residual connections (~300K params) enables real-time CPU inference < 10 ms",
    "Three-zone safety layer prevents the majority of collisions near obstacles",
    "Recovery FSM handles deadlock situations autonomously",
    "Fully dockerized pipeline enables reproducible evaluation across all 360 worlds",
    "Headless video recording captures simulation results for analysis",
]
for c in conclusions:
    add_para(tf, f"• {c}", font_size=13, color=LIGHT_GRAY, space_before=Pt(6))

# Right - future work
tf2 = add_text(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.4),
               "🔮  Future Work", font_size=20, color=ORANGE, bold=True)
future = [
    "DAgger/online fine-tuning to correct distribution shift errors",
    "Transformer-based model (~600K params) for richer temporal context",
    "Curriculum learning: easy→hard world scheduling",
    "Full DynaBARN integration with 2-frame LiDAR velocity estimation",
    "Sim-to-real transfer for physical Jackal deployment at ICRA 2026",
    "Run complete 360-world evaluation for competition submission",
]
for f in future:
    add_para(tf2, f"• {f}", font_size=13, color=LIGHT_GRAY, space_before=Pt(6))

# Bottom box - limitations
box = add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), BG_MID)
tf3 = add_text(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.4),
               "⚠️  Current Limitations", font_size=16, color=GOLD, bold=True)
add_para(tf3, "Limited real-world testing  •  World 0 shows collision (need more training data / fine-tuning)  •  "
              "DynaBARN layer not yet deployed  •  No GPU acceleration available on target hardware",
         font_size=12, color=LIGHT_GRAY, space_before=Pt(6))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — TEAM CONTRIBUTIONS
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Team Contributions", "Team 14 — IIT Kharagpur")

team = [
    ("Pampana Charan Teja", "25AI60R22", "Leader", [
        "BC Model design & architecture",
        "PyTorch training pipeline",
        "ONNX model export & optimization",
        "Overall project coordination",
    ], ACCENT2),
    ("Pucha Arun Kumar", "25AI60R26", "Data Engineer", [
        "Expert data collection (KUL+FM planner)",
        "Dataset validation & quality checks",
        "Data pipeline & preprocessing",
    ], GREEN),
    ("Arkoti Charan Teja", "25AI60R19", "Systems Engineer", [
        "Safety layer implementation",
        "Recovery FSM design & tuning",
        "ROS node integration (nav_node.py)",
    ], ORANGE),
    ("Bommireddi Nagendra", "25AI60R14", "Research Engineer", [
        "DynaBARN obstacle handling layer",
        "Data augmentation strategies",
        "Ablation studies & hyperparameter tuning",
    ], RGBColor(0xA7, 0x8B, 0xFA)),
    ("Adarsh Raj", "25AI560R25", "DevOps Engineer", [
        "Docker & Singularity container setup",
        "Evaluation benchmarking scripts",
        "Batch runner & video recording pipeline",
    ], RGBColor(0xF4, 0x72, 0xB6)),
]

for i, (name, roll, role, contribs, color) in enumerate(team):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.2)
    y = Inches(1.3 + row * 3.1)
    w = Inches(3.9)
    h = Inches(2.8 if row == 0 else 2.6)

    card = add_rounded_rect(slide, x, y, w, h, BG_MID)
    # Color strip
    add_shape(slide, x, y, Inches(0.08), h, color)
    # Name and roll
    add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(3.4), Inches(0.4),
             name, font_size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.55), Inches(3.4), Inches(0.3),
             f"{roll}  •  {role}", font_size=12, color=color, bold=True)
    # Contributions
    tf = add_text(slide, x + Inches(0.25), y + Inches(0.95), Inches(3.4), Inches(1.6),
                  contribs[0], font_size=12, color=LIGHT_GRAY)
    for c in contribs[1:]:
        add_para(tf, f"• {c}", font_size=12, color=LIGHT_GRAY, space_before=Pt(3))

# References at bottom
add_text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
         "References:  BARN Challenge 2026  •  LiCS-KI (arXiv:2406.14947)  •  github.com/Daffan/the-barn-challenge",
         font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════════
os.makedirs(os.path.dirname(OUTPUT_FILE), exist_ok=True)
prs.save(OUTPUT_FILE)
print(f"\n🎉 Presentation saved to: {OUTPUT_FILE}")
print(f"   Total slides: {len(prs.slides)}")
